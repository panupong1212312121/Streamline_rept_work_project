from img2table.document import Image
import tabula
import fitz as pyMuPDF
from docx import Document
from pptx import Presentation

import pandas as pd

class ExtractTable:
    
    def image(self,file_afterTransfrom):
        dfs = Image(file_afterTransfrom)
        return dfs
    
    def pdf(self,file,file_afterTransfrom):
        dfs = {}
        page_amount = pyMuPDF.open(file_afterTransfrom).page_count
        for page_num in range(page_amount):
            tables = tabula.read_pdf(file, pages=page_num+1)
            if tables:
                dfs[page_num+1] = [table for table in tables]
        return dfs
    
    def docx(self,file_afterTransfrom):
        document = Document(file_afterTransfrom)
        dfs = []
        for table in document.tables:
            df = [['' for _ in range(len(table.columns))] for _ in range(len(table.rows))]
            for i, row in enumerate(table.rows):
                for j, cell in enumerate(row.cells):
                    if cell.text:
                        df[i][j] = cell.text
            dfs.append(pd.DataFrame(df))
        return dfs
    
    def pptx(self,file_afterTransfrom):
        presentation = Presentation(file_afterTransfrom)
        dfs = {}
        page = 0
        for slide in presentation.slides:
            df = []
            page += 1
            hasTable = True
            for shape in slide.shapes:
                if not shape.has_table:
                    hasTable = False
                    continue
                table = shape.table
                row_count = len(table.rows)
                col_count = len(table.columns)
                table_data = []
                for row in range(row_count):
                    row_values = []
                    for col in range(col_count):
                        cell = table.cell(row, col)
                        text = ""
                        for paragraph in cell.text_frame.paragraphs:
                            for run in paragraph.runs:
                                text += run.text
                        row_values.append(text)
                    table_data.append(row_values)
                df.append(pd.DataFrame(table_data))
            if hasTable:
                dfs[page] = df
        return dfs
    
    def link(self,url_name):
        df = pd.read_html(url_name)
        return df 
             